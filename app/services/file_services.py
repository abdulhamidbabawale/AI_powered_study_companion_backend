import fitz  # pymupdf
from pptx import Presentation

def extract_text_from_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

def extract_text_from_pptx(file_bytes: bytes) -> str:
    import io
    prs = Presentation(io.BytesIO(file_bytes))
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text.append(shape.text_frame.text)
    return "\n".join(text)